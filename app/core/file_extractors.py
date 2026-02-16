import pandas as pd
from pathlib import Path
from pdfminer.high_level import extract_text as pdf_text
from docx import Document as Docx
from pptx import Presentation


def _extract_excel(path: Path) -> str:
    try:
        if path.suffix.lower() == ".csv":
            df = pd.read_csv(path, dtype=str, keep_default_na=False)
            return "\n".join(df.astype(str).agg(" ".join, axis=1))

        sheets = pd.read_excel(
            path,
            sheet_name=None,
            dtype=str,
            engine="openpyxl"
        )
        rows = []
        for df in sheets.values():
            rows += df.astype(str).agg(" ".join, axis=1).tolist()

        return "\n".join(rows)
    except Exception:
        return ""


def extract_text(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    try:
        if ext in {"txt", "py", "java"}:
            return path.read_text("utf-8", errors="ignore")

        if ext in {"csv", "xls", "xlsx"}:
            return _extract_excel(path)

        if ext == "docx":
            return "\n".join(p.text for p in Docx(path).paragraphs)

        if ext == "pptx":
            return "\n".join(
                sh.text
                for sl in Presentation(path).slides
                for sh in sl.shapes
                if hasattr(sh, "text")
            )

        if ext == "pdf":
            return pdf_text(path)

    except Exception:
        return ""

    return ""
